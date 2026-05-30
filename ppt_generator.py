#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""依資安事件資料自動產生 PowerPoint 分析簡報（python-pptx）"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ── 調色盤 ──
NAVY = RGBColor(0x14, 0x21, 0x3D)
NAVY2 = RGBColor(0x1E, 0x2D, 0x4F)
TEAL = RGBColor(0x2E, 0xC4, 0xB6)
RED = RGBColor(0xE6, 0x39, 0x46)
AMBER = RGBColor(0xF4, 0xA2, 0x61)
CYAN = RGBColor(0x27, 0xA6, 0xC4)
PURPLE = RGBColor(0x8A, 0x6F, 0xB0)
LIGHT = RGBColor(0xF6, 0xF8, 0xFC)
ROW = RGBColor(0xEE, 0xF2, 0xF9)
DARK = RGBColor(0x1A, 0x22, 0x33)
GRAY = RGBColor(0x5C, 0x66, 0x79)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ICE = RGBColor(0xCA, 0xDC, 0xFC)

FONT = "Microsoft JhengHei"

TYPE_COLOR = {
    "駭客攻擊": RED,
    "勒索病毒/加密": AMBER,
    "個資外洩": CYAN,
    "其他": GRAY,
}


def _set_run(run, text, size, color, bold=False, font=FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font)


def _text(slide, text, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT,
          anchor=MSO_ANCHOR.TOP, line_spacing=None):
    """加入單段文字方塊"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    _set_run(p.add_run(), text, size, color, bold)
    return box


def _rich(slide, runs, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=None):
    """加入多段(rich)文字：runs = [(text,size,color,bold), ...]"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    for (t, sz, col, bd) in runs:
        _set_run(p.add_run(), t, sz, col, bd)
    return box


def _rect(slide, x, y, w, h, fill, shape=MSO_SHAPE.RECTANGLE, line=None, shadow=False):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    # shadow=True 的卡片改用細邊框增加層次（自訂陰影 XML 會被 PowerPoint 視為損毀）
    if shadow:
        sp.line.color.rgb = RGBColor(0xDD, 0xE3, 0xEE)
        sp.line.width = Pt(1)
    elif line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def _truncate(s, n):
    s = (s or "").strip()
    return s if len(s) <= n else s[:n] + "…"


def _slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 空白版面


def _bg(slide, color, w, h):
    _rect(slide, 0, 0, w, h, color)


def build_pptx(announcements, out_path, report_date, period_start):
    """產生分析簡報。announcements 需含 code/company/date/title/event_type/cause/impact/future"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    W, H = 13.333, 7.5

    # 分類統計
    counts = {}
    for a in announcements:
        counts[a.get("event_type", "其他")] = counts.get(a.get("event_type", "其他"), 0) + 1

    # ===== Slide 1：封面 =====
    s = _slide(prs)
    _bg(s, NAVY, W, H)
    _rect(s, 11.3, -1.2, 3.6, 3.6, NAVY2, shape=MSO_SHAPE.OVAL)
    _rect(s, 12.2, 5.2, 2.8, 2.8, NAVY2, shape=MSO_SHAPE.OVAL)
    _rect(s, 0.7, 0.95, 0.5, 0.5, TEAL, shape=MSO_SHAPE.OVAL)
    _text(s, "資安重訊分析", 1.4, 0.95, 5, 0.5, 14, TEAL, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    _text(s, "上市櫃公司\n資安事件因應分析", 0.7, 2.0, 11.5, 2.2, 44, WHITE, bold=True, line_spacing=54)
    _text(s, f"資料期間 {period_start} ~ {report_date}　｜　MOPS 重大訊息全文檢索　｜　{len(announcements)} 家公司",
          0.72, 4.5, 11.5, 0.5, 16, ICE)
    _rect(s, 0.75, 5.15, 3.2, 0.04, TEAL)
    _text(s, "資料來源：台灣證券交易所公開資訊觀測站（mopsov.twse.com.tw）", 0.72, 6.7, 11, 0.4, 11, RGBColor(0x88, 0x93, 0xAD))

    # ===== Slide 2：事件總覽表 =====
    s = _slide(prs)
    _bg(s, LIGHT, W, H)
    _text(s, "事件總覽", 0.6, 0.4, 8, 0.7, 30, DARK, bold=True)
    _text(s, f"資料期間內 MOPS 揭露之 {len(announcements)} 起資安相關重大訊息", 0.62, 1.1, 10, 0.4, 13, GRAY)

    stat_items = [(str(len(announcements)), "起資安事件", TEAL),
                  (str(len(counts)), "大事件類型", CYAN),
                  ("第26款", "資通安全事件", RED)]
    for i, (big, lab, col) in enumerate(stat_items):
        x = 0.6 + i * 2.4
        _rect(s, x, 1.6, 2.15, 1.05, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        _text(s, big, x, 1.66, 2.15, 0.6, 27, col, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, lab, x, 2.22, 2.15, 0.36, 11.5, GRAY, align=PP_ALIGN.CENTER)

    # 表格
    headers = ["公司", "代號", "事件日", "事件類型", "公司揭露之影響"]
    col_w = [2.0, 0.95, 1.05, 2.5, 5.6]
    n = len(announcements) + 1
    tbl_shape = s.shapes.add_table(n, len(headers), Inches(0.6), Inches(2.95),
                                   Inches(sum(col_w)), Inches(min(3.9, 0.42 * n)))
    table = tbl_shape.table
    table.first_row = False
    table.horz_banding = False
    for j, cw in enumerate(col_w):
        table.columns[j].width = Inches(cw)
    for j, htext in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.08); cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
        para = cell.text_frame.paragraphs[0]
        _set_run(para.add_run(), htext, 12.5, WHITE, bold=True)
    for i, a in enumerate(announcements, start=1):
        impact = _truncate(a.get("impact", "") or "—", 34)
        date_short = a["date"].split("/", 1)[1] if a["date"].count("/") >= 2 else a["date"]
        vals = [a["company"], a["code"], date_short,
                a.get("event_type", "其他"), impact]
        for j, v in enumerate(vals):
            cell = table.cell(i, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if i % 2 else ROW
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08); cell.margin_top = Inches(0.01); cell.margin_bottom = Inches(0.01)
            para = cell.text_frame.paragraphs[0]
            col = TYPE_COLOR.get(v, DARK) if j == 3 else DARK
            _set_run(para.add_run(), v, 11.5, col, bold=(j == 3))

    # ===== Slide 3：類型分布 =====
    s = _slide(prs)
    _bg(s, LIGHT, W, H)
    _text(s, "事件類型分布", 0.6, 0.4, 9, 0.7, 30, DARK, bold=True)
    _text(s, "依公告主旨與內文關鍵字分類", 0.62, 1.1, 10, 0.4, 13, GRAY)

    order = ["駭客攻擊", "勒索病毒/加密", "個資外洩", "其他"]
    present = [(t, counts[t]) for t in order if counts.get(t)]
    chart_data = CategoryChartData()
    chart_data.categories = [t for t, _ in present]
    chart_data.add_series("事件數", [c for _, c in present])
    gframe = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(0.5), Inches(1.7),
                                Inches(5.2), Inches(4.9), chart_data)
    chart = gframe.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(12)
    chart.legend.font.name = FONT
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = "0"
    plot.data_labels.number_format_is_linked = False
    plot.data_labels.font.size = Pt(13)
    plot.data_labels.font.bold = True
    plot.data_labels.font.color.rgb = WHITE
    pts = plot.series[0].points
    for idx, (t, _) in enumerate(present):
        pts[idx].format.fill.solid()
        pts[idx].format.fill.fore_color.rgb = TYPE_COLOR.get(t, GRAY)

    # 右側卡片
    desc = {
        "駭客攻擊": "系統遭外部入侵或企圖登入，立即啟動防禦機制",
        "勒索病毒/加密": "伺服器資料遭加密，須系統隔離與資料復原",
        "個資外洩": "源自開發流程或委外廠商，須通報並通知當事人",
        "其他": "其他資安相關事件",
    }
    names_by_type = {}
    for a in announcements:
        names_by_type.setdefault(a.get("event_type", "其他"), []).append(a["company"])
    for i, (t, c) in enumerate(present):
        y = 1.85 + i * 1.45
        col = TYPE_COLOR.get(t, GRAY)
        _rect(s, 6.2, y, 6.6, 1.28, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        _rect(s, 6.2, y, 0.12, 1.28, col)
        _text(s, t, 6.45, y + 0.12, 4.6, 0.4, 16, DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, f"{c} 起", 11.0, y + 0.12, 1.6, 0.4, 18, col, bold=True, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, "、".join(names_by_type.get(t, []))[:40], 6.45, y + 0.55, 6.1, 0.3, 11.5, DARK, bold=True)
        _text(s, desc.get(t, ""), 6.45, y + 0.84, 6.2, 0.36, 11, GRAY)

    # ===== Slide 4：共通因應流程 =====
    s = _slide(prs)
    _bg(s, LIGHT, W, H)
    _text(s, "共通因應流程", 0.6, 0.4, 9, 0.7, 30, DARK, bold=True)
    _text(s, "各公司處理過程高度一致，已形成標準應變劇本", 0.62, 1.1, 11, 0.4, 13, GRAY)
    steps = [("01", "偵測異常", "資安單位或外部通報\n發現系統異常或入侵"),
             ("02", "啟動應變", "立即啟動資安應變機制\n系統隔離、防護強化"),
             ("03", "外部協助", "委請外部資安專業團隊\n鑑識調查與資料復原"),
             ("04", "強化監控", "盤查系統、提升架構\n持續密切監控防護")]
    for i, (num, title, body) in enumerate(steps):
        x = 0.7 + i * 3.07
        _rect(s, x, 2.0, 2.75, 2.5, NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        _rect(s, x + 1.02, 2.25, 0.72, 0.72, TEAL, shape=MSO_SHAPE.OVAL)
        _text(s, num, x + 1.02, 2.25, 0.72, 0.72, 19, NAVY, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, title, x, 3.15, 2.75, 0.5, 18, WHITE, bold=True, align=PP_ALIGN.CENTER)
        _text(s, body, x + 0.15, 3.68, 2.45, 0.75, 11.5, ICE, align=PP_ALIGN.CENTER, line_spacing=15)
        if i < 3:
            _text(s, "›", x + 2.72, 2.0, 0.4, 2.5, 28, GRAY, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _rect(s, 0.7, 5.0, 11.9, 1.45, RGBColor(0xFF, 0xF1, 0xF0), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _rect(s, 0.7, 5.0, 0.12, 1.45, RED)
    _rich(s, [("個資外洩事件額外步驟　", 15, RED, True), ("（依事件性質適用）", 12, GRAY, False)],
          0.95, 5.15, 11.4, 0.45, anchor=MSO_ANCHOR.MIDDLE)
    _text(s, "依法通報主管機關　•　主動通知受影響當事人並提供客服　•　向調查局報案　•　強化委外廠商與第三方工具管理",
          0.95, 5.65, 11.5, 0.6, 12, DARK, line_spacing=16)

    # ===== Slide 5+：個案內文卡片（每頁 4 張）=====
    per = 4
    for page in range((len(announcements) + per - 1) // per):
        s = _slide(prs)
        _bg(s, LIGHT, W, H)
        suffix = "" if len(announcements) <= per else f"（{page + 1}）"
        _text(s, f"個案因應重點{suffix}", 0.6, 0.4, 10, 0.7, 30, DARK, bold=True)
        _text(s, "各公司公告之發生緣由與因應措施摘要", 0.62, 1.1, 11, 0.4, 13, GRAY)
        chunk = announcements[page * per:(page + 1) * per]
        for k, a in enumerate(chunk):
            col_i = k % 2
            row_i = k // 2
            x = 0.6 + col_i * 6.15
            y = 1.7 + row_i * 2.55
            col = TYPE_COLOR.get(a.get("event_type", "其他"), GRAY)
            _rect(s, x, y, 5.9, 2.35, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
            _rect(s, x, y, 5.9, 0.12, col)
            _rich(s, [(f"{a['company']} ", 16, DARK, True), (f"{a['code']}", 13, GRAY, False)],
                  x + 0.25, y + 0.24, 4.0, 0.4, anchor=MSO_ANCHOR.MIDDLE)
            # 類型徽章
            badge_w = 1.35
            _rect(s, x + 5.9 - badge_w - 0.2, y + 0.26, badge_w, 0.38, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            _text(s, a.get("event_type", "其他"), x + 5.9 - badge_w - 0.2, y + 0.26, badge_w, 0.38,
                  10.5, WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _rich(s, [("緣由　", 11, col, True), (_truncate(a.get("cause", ""), 48), 11, DARK, False)],
                  x + 0.25, y + 0.78, 5.4, 0.7, line_spacing=15)
            _rich(s, [("因應　", 11, col, True), (_truncate(a.get("future", "") or a.get("impact", ""), 60), 11, DARK, False)],
                  x + 0.25, y + 1.5, 5.4, 0.78, line_spacing=15)

    prs.save(out_path)
    return out_path
